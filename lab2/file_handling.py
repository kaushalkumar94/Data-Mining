#!/usr/bin/env python
# coding: utf-8

# In[2]:


get_ipython().system(' pip install pandas')


# ## CSV file

# In[35]:


# Import the pandas library
import pandas as pd

# Define the path of the CSV file
# Use a raw string (r"") to avoid issues with backslashes in Windows paths
file_path = r"D:\CODING\data mining\lab 2\dataset\addresses.csv"

# Read the CSV file into a Pandas DataFrame
df = pd.read_csv(file_path)

# Display the DataFrame
print(df)


# In[5]:


print(df.head())


# In[6]:


print(df.info())


# In[7]:


print(df.describe())


# ## EXCEL file

# In[8]:


# Import the pandas library
import pandas as pd

# Define the path of the Excel file
# Raw string is used to correctly handle Windows file paths
file_path = r"D:\CODING\data mining\lab 2\dataset\file_example_XLSX_1000.xlsx"

# Read the Excel file into a Pandas DataFrame
# By default, the first sheet is read
df_excel = pd.read_excel(file_path)

# Display the DataFrame
print(df_excel)


# In[10]:


print(df_excel.head())


# In[11]:


print(df_excel.info())


# In[12]:


print(df_excel.describe())


# ## TEXT file 

# In[14]:


# Import the pandas library
import pandas as pd

# Define the path of the text file
# Raw string is used to correctly handle Windows file paths
file_path = r"D:\CODING\data mining\lab 2\dataset\frankenstein.txt"

# Open the text file using Python's built-in file handling
with open(file_path, "r", encoding="utf-8") as file:
    # Read all lines from the text file
    lines = file.read().splitlines()

# Convert the text data into a Pandas DataFrame
# Each line of the text file becomes a separate row
df_txt = pd.DataFrame(lines, columns=["Text"])

# Display the DataFrame
print(df_txt)


# In[15]:


print(df_txt.head(10))


# In[16]:


print("Total lines:", len(df_txt))


# ## JSON file

# In[24]:


# Import the pandas library
import pandas as pd

# Define the path of the JSON file
# Raw string is used to correctly handle Windows file paths
file_path = r"D:\CODING\data mining\lab 2\dataset\price.json"

# Read the JSON file into a Pandas DataFrame
# Works best when JSON is in list-of-records format
df_json = pd.read_json(file_path)

# Display the DataFrame
print(df_json)


# In[26]:


print(df_json.head(10))


# In[30]:


print(df_json.describe())


# In[33]:


print(df_json.info())


# ## Opening data by url on python 

# In[39]:


get_ipython().system(' pip install lxml')


# In[ ]:


# Import required libraries
import pandas as pd
import requests

# Wikipedia URL
url = "https://en.wikipedia.org/wiki/List_of_countries_and_dependencies_by_population"

# Send request with browser-like User-Agent
response = requests.get(
    url,
    headers={"User-Agent": "Mozilla/5.0"}
)

# Read HTML tables from the page content
tables = pd.read_html(response.text)

# Select the first table
df_wiki = tables[0]

# Display data in tabular form
print(df_wiki)


# ## Read pdf files

# In[51]:


get_ipython().system('pip install pdfplumber')


# In[52]:


import pdfplumber
import pandas as pd

# PDF file path (can be local or URL-downloaded)
file_path = "D:\STUDY\GATE\PYQS\CS\CS22025.pdf"

# List to store extracted text
data = []

# Open the PDF file
with pdfplumber.open(file_path) as pdf:
    for page in pdf.pages:
        text = page.extract_text()
        if text:
            for line in text.split("\n"):
                data.append([line])

# Convert extracted text into DataFrame
df_pdf = pd.DataFrame(data, columns=["PDF_Text"])

# Display PDF data as table
print(df_pdf)


# ## Read PPTx

# In[56]:


get_ipython().system('pip install python-pptx')


# In[60]:


from pptx import Presentation
import pandas as pd

# PPT file path
file_path = "D:\\CODING\\data mining\\lab 2\\dataset\\file_example_PPT_250kB.pptx"

# Load the PowerPoint file
presentation = Presentation(file_path)

# List to store slide text
data = []

# Extract text from each slide
for slide_number, slide in enumerate(presentation.slides, start=1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            data.append([slide_number, shape.text])

# Convert extracted text into DataFrame
df_ppt = pd.DataFrame(data, columns=["Slide_Number", "Slide_Text"])

# Display PPT data as table
print(df_ppt)


# ## Opening text file without pyhton

# In[66]:


with open(r"D:\CODING\data mining\lab 2\dataset\frankenstein.txt", 
          "r", encoding="utf-8") as file:
    content = file.read()

print(content[:500])


# ## Opening Multimedia file such as image , video , audio

# In[68]:


get_ipython().system(' pip install pillow')


# In[70]:


from PIL import Image

# Use raw string to avoid path errors
file_path = r"D:\CODING\data mining\lab 2\dataset\Screenshot 2026-02-01 150322.png"

# Open the image
img = Image.open(file_path)

# Display the image
img.show()


# In[ ]:


get_ipython().system('pip install opencv-python')


# In[76]:


import cv2

# Correct image path
file_path = r"D:\CODING\data mining\lab 2\dataset\Screenshot 2026-02-01 150322.png"

# Read image
img = cv2.imread(file_path)

# Display image in a window
cv2.imshow("Image", img)
cv2.waitKey(0)
cv2.destroyAllWindows()


# ## Reading data from database
# 

# In[ ]:


pip install mysql-connector-python pandas sqlalchemy


# In[12]:


import pandas as pd
from sqlalchemy import create_engine

# MySQL connection details (from MySQL Workbench)
username = "root"
password = "11122004@kK"
host = "127.0.0.1"
port = "3306"
database = "parks_and_recreation"

# Create MySQL connection engine
engine = create_engine(
    f"mysql+mysqlconnector://{username}:{password}@{host}:{port}/{database}"
)

# Example query (change table name if needed)
query = "SELECT * FROM employee_demographics"

# Read data from MySQL into DataFrame
df = pd.read_sql(query, engine)

# Display data in tabular form
print(df)


# Unknown MySQL Server Host Error
# 
# Problem:
# MySQL reported unknown host during connection.
# 
# Reason:
# Special character @ in the password was misinterpreted in the connection string.
# 
# Solution:
# URL-encoded the password or used mysql.connector directly confirming credentials.

# In[11]:


import pandas as pd
from sqlalchemy import create_engine

# MySQL connection details
username = "root"
password = "11122004%40kK"   # @ replaced with %40
host = "127.0.0.1"
port = "3306"
database = "parks_and_recreation"

# Create MySQL engine
engine = create_engine(
    f"mysql+mysqlconnector://{username}:{password}@{host}:{port}/{database}"
)

# Test query (replace table name if needed)
query = "SELECT * FROM employee_demographics"

# Read data into DataFrame
df = pd.read_sql(query, engine)

# Display output
print(df)


# In[10]:


import mysql.connector
import pandas as pd

conn = mysql.connector.connect(
    host="localhost",
    user="root",
    password="11122004@kK",
    database="parks_and_recreation",
    auth_plugin="mysql_native_password"
)

df = pd.read_sql("SELECT * FROM employee_demographics", conn)
print(df)

conn.close()


# In[9]:


# Import required libraries
import mysql.connector
import pandas as pd

# Establish connection to MySQL database
conn = mysql.connector.connect(
    host="localhost",
    user="python_user",
    password="python123",
    database="parks_and_recreation"
)

# SQL query to fetch data from a table
# Change table name if required
query = "SELECT * FROM employee_demographics"

# Read data from MySQL into a Pandas DataFrame
df = pd.read_sql(query, conn)

# Display the dataset in tabular form
print(df)

# Close the database connection
conn.close()

